#!/usr/bin/env python3
"""
Build TravelTrust Pitch Deck PPTX (and optional PDF) from the Release storyboard — **15 slides** (compressed main deck: vector travel loop, protocol stack; **Why now** immediately after vectors; three lines → one order → strengths → **FeeRouter** disclosure boundary; roadmap+disclosure rail; competitive matrix; journey; closing routes DD to FAQ/Whitepaper). Target **~7–10 min** live; long copy lives in **speaker notes** / appendix.

Requires: pip install python-pptx
PDF: LibreOffice (soffice) headless from PPTX (required for handoff; no fpdf fallback).

Design system: `ir_deck_design_system.py`. Storytelling IA (Problem→WhyNow→WhyUs→WhyWin→DD→Close): `ir_deck_storytelling.py`
(speaker notes + left rail + eyebrows; **no** new economic numbers). Optional slide audit: set **`TRAVELTRUST_IR_SLIDE_AUDIT=1`** to write `docs/fundraising/internal/37-IR-Pitch-Deck-Slide-Audit-RUN.md` (maintainer QA; **not** a default ship artifact).
Optional dark theme: TRAVELTRUST_DECK_THEME=dark. Spec: docs/fundraising/internal/36-IR-Pitch-Deck-Design-System.md.

Writes to docs/fundraising/external/export-ready/:
  04-PitchDeck-v{release}-CN.pdf and 04-PitchDeck-v{release}-EN.pdf only (slot 04).

PPTX → docs/fundraising/internal/deck-editable/. IC appendix is internal-only (`TRAVELTRUST_BUILD_IC_INTERNAL=1`).
"""
from __future__ import annotations

import json
import os
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
ANCHORS = ROOT / "registry" / "fundraising-external-numeric-anchors.v1.json"
OUT_DIR = ROOT / "docs" / "fundraising" / "external" / "export-ready"

_TOOLS = Path(__file__).resolve().parent
if str(_TOOLS) not in sys.path:
    sys.path.insert(0, str(_TOOLS))
import investor_handoff_layout as ihr  # noqa: E402
import ir_deck_design_system as ids  # noqa: E402
import ir_deck_slide_audit as ids_audit  # noqa: E402
import ir_deck_storytelling as ids_story  # noqa: E402
from investor_ir_office import convert_pptx_to_pdf, find_soffice  # noqa: E402


SlideSpec = tuple[str, list[str], list[str]]


def _unpack_slide(spec: tuple[str, ...] | SlideSpec) -> tuple[str, list[str], list[str]]:
    if len(spec) == 2:
        title, bullets = spec  # type: ignore[misc]
        return str(title), list(bullets), []
    title, vis, notes = spec  # type: ignore[misc]
    return str(title), list(vis), list(notes)


def _slides_audit_pairs(slides: list[tuple[str, ...] | SlideSpec]) -> list[tuple[str, list[str]]]:
    """Visible copy only — matches on-slide body for read-time heuristics."""
    return [(t, v) for (t, v, _) in (_unpack_slide(s) for s in slides)]


def _split_competitive_columns(bullets: list[str]) -> tuple[tuple[str, str], tuple[str, str], tuple[str, str], str] | None:
    if len(bullets) < 4:
        return None

    def one(line: str) -> tuple[str, str]:
        for sep in ("：", ":"):
            if sep in line:
                a, b = line.split(sep, 1)
                return a.strip(), b.strip()
        return line.strip(), ""

    return one(bullets[0]), one(bullets[1]), one(bullets[2]), bullets[3]


def _pillar_parts(title: str, bullets: list[str]) -> list[str] | None:
    if ("优势" in title or "strengths" in title.lower()) and len(bullets) == 1 and "·" in bullets[0]:
        s = bullets[0].replace(" · ", "·").replace(" • ", "·")
        parts = [p.strip() for p in s.split("·") if p.strip()]
        return parts or None
    return None


def _slide_layout_kind(title: str, bullets: list[str]) -> str:
    t = title.lower()
    if (
        "单笔订单对齐交付与争议" in title
        or "trip-order: align delivery and disputes" in t
        or "主链" in title
        or "deal loop" in t
    ):
        return "vector_loop"
    if (
        "订单与规则分几层" in title
        or "钱和规则分几层" in title
        or "trip-order & rules stack" in t
        or "money & rules stack" in t
        or "协议栈" in title
        or "protocol stack" in t
        or "order rules |" in t
    ):
        return "vector_stack"
    if (
        "feerouter" in t
        or "协议经济·第一层" in title
        or ("protocol economics" in t and "layer 1" in t)
        or "disclosure boundary" in t
        or "披露边界" in title
    ):
        return "vector_fee"
    if "竞争" in title or "competitive" in t:
        return "matrix"
    if "dd" in t or "q&a" in t or "问答" in title:
        return "faq"
    if "痛点" in title or t == "pain":
        return "pain_grid"
    if "路线图" in title or "roadmap" in t:
        return "roadmap"
    if ("一句话" in title or "one line" in t) and len(bullets) <= 1:
        return "hero"
    if _pillar_parts(title, bullets):
        return "pillars"
    if len(bullets) >= 4:
        return "cards_2col"
    if len(bullets) >= 3:
        return "cards_2col"
    return "linear"


def _fee_slide_chrome_split(title: str, lang: str) -> tuple[str, str] | None:
    """Long-title safe chrome: disclosure lead line + parenthetical tail (FeeRouter unchanged in tail).

    Uses full-width ``（`` for CN and the first ``(`` for EN so speaker-note / audit titles stay verbatim.
    """
    sep = "（" if lang == "CN" else "("
    if sep not in title:
        return None
    i = title.index(sep)
    lead = title[:i].strip()
    tail = title[i:].strip()
    if not lead or not tail:
        return None
    return lead, tail


def _pitch_chrome_long_split(title: str, lang: str, slide_idx: int) -> tuple[str, str] | None:
    """Chrome-only two-line split for slides where a single wrapped title clips the body band.

    Speaker notes still use the full ``title`` from the storyboard tuple. Targets are curated
    by (slide_idx, lang); split on full-width ``（`` (CN) or the first ``(`` (EN) like FeeRouter slides.
    """
    if (slide_idx, lang) not in {(5, "CN"), (5, "EN"), (9, "EN")}:
        return None
    sep = "（" if lang == "CN" else "("
    if sep not in title:
        return None
    i = title.index(sep)
    lead = title[:i].strip()
    tail = title[i:].strip()
    if not lead or not tail:
        return None
    return lead, tail


SLIDES_CN: list[SlideSpec] = [
    (
        "TravelTrust",
        [
            "旅行订单金融科技",
            "托管 · 履约 · 争议同链",
            "以已发布公告与白皮书为边界",
            "误读防火墙：单笔订单信托责+托管/争议；非旅费支付、非大盘 OTA、非代币募资页（口径同 06 / FAQ）。",
        ],
        [],
    ),
    (
        "TravelTrust",
        [
            "旅行订单金融科技：围绕单笔订单对齐交付与争议。",
            "边界：旅费本金与平台经常性收入分列披露；定量与条款以白皮书 + NDA 为准。",
            "跟投锚（定性）：分销碎片化下，单笔订单级「资金状态 + 履约/争议」可对账、可审计的硬需求；以下页为同一单笔主链递进。",
            "Why us（定性）：旅资状态机·服务费·治理三轨与 06 / FAQ 同句可读；非旅费支付中台叙事。",
        ],
        [
            "定位（仅 Notes）：平台体验 + 可嵌入托管/争议基建；不争「一夜替代存量 OTA」—与第 9 页总括同向，不在此屏堆叠。",
        ],
    ),
    (
        "核心痛点",
        [
            "信息对不上",
            "争议难审计、对账成本高",
            "订单资金状态：客人侧与平台侧是否两本账？",
        ],
        [],
    ),
    (
        "用户履约闭环｜单笔订单：交付与争议对齐（同一单笔｜不讲费项拆分）",
        [],
        [
            "讲者首句：这一页只回答一单在履约侧怎么走通；不讲池子与比例。",
            "图示为单笔交易闭环；勿补充材料未载明的数字。",
            "屏上矢量仅为心智模型；口播开头先念：当前演示/测试网或生产以已发布公告与环境标签为准（与第 7 页一致）。",
            "顺序：找到→下单→对齐交付→履约/争议→以已发布公告与披露为准。",
            "节奏与覆盖：以已发布公告与披露为准；勿暗示未公告范围、勿超出已披露表述。",
        ],
    ),
    (
        "控制面·协议栈与订单规则｜责任切片（同一单笔｜规则/披露/环境分层；非第二条用户动线）",
        [],
        [
            "讲者首句：仍是同一单笔；这一页回答控制面怎么分层——不是另一条业务流程。",
            "与上页同一单笔；纵向栈示意责任切片（控制层）；勿引入材料未载明的量化表述。",
            "屏上五步栈若带缩写图例：口播先念环境标签全称，再解释层级；禁止暗示未公告已全量生产。",
            "五步顺序不变：找到→下单→对齐交付→履约/争议→以已发布公告与披露为准。",
            "节奏与覆盖：以已发布公告与披露为准；勿暗示未公告范围。命名与分配细节会后沟通；未公告前勿断言全面上线。",
        ],
    ),
    (
        "为何现在是现在（定性）",
        [
            "需求：客人要听得真，不要吹。（结构见 06 §2.4）",
            "供给：名声带走；披露随已发布公告。",
            "监管：订单资金状态与平台服务费分开走。（口径见 06）",
            "窗口：分销碎，先对齐交付再接。（同节窗口论述）可复读：06 已发布节号 + 公告日期 + 环境标签。",
        ],
        [
            "Timing 复读锚：窗口叙事与 06 已发布节号、公告日期、环境标签对读；行业比例表在 NDA。",
            "行业对比与数据表多在保密前提下讨论；本段不作认购推介。",
        ],
    ),
    (
        "三句话记住",
        [
            "平台赚服务费、抽成；旅费不算平台自有收入",
            "治理信号：只投票；非旅费支付工具。",
            "披露节奏：已发布公告与版本/环境标签（细则见白皮书）·壁垒侧重披露与对账可复读",
        ],
        [
            "第一行：经常性收入来自服务费与抽成；旅费本金不计入平台自有收入。",
            "第二、三行见上屏字面；不在此页展开费项拆分。",
            "三条可独立复述；进一步拆分留至会后。",
        ],
    ),
    (
        "一笔订单：交付对齐了吗",
        [
            "一单从下单一路对齐交付？",
            "进度与边界：对照已发布公告与环境标签。",
            "材料：同包 FAQ / 白皮书；口头不合并 45/55 与其他费用分母。",
        ],
        [
            "盯一单：从下单到对齐交付；讲到哪一步，勿超过已发布公告与环境标签所写范围。",
        ],
    ),
    (
        "系统性优势（投资视角·定性：信任栈可嵌入渠道；壁垒在披露与对账闭环）",
        [
            "履约侧 · 订单对齐 · 交付与争议 · 出事收口 · 随已发布公告/标签",
        ],
        [
            "总括（仅 Notes）：平台体验 + 可嵌入托管/争议基建；不争「一夜替代存量」。",
            "上屏为分项柱条，与第 4–5 页矢量主链同向复述。",
        ],
    ),
    (
        "已发布口径·披露边界（外环披露｜FeeRouter｜屏上仅已发布口径）",
        [],
        [
            "讲者首句：这是披露边界第一层示意（外环披露）；先对齐分母，不讲营业账。",
            "披露口径边界示意；屏上只印 100/45/55。",
            "口播三句分母（与 FAQ / 白皮书一致）：披露·45% 与披露·55% 为可分配手续费第一层拆分；链上油费/仲裁/罚没不在同一分母。",
            "披露·45%：一国一桶；桶里再切在网签。",
            "披露·55%：Global Pool；再切在网签。",
            "链上油费/仲裁不是这盘账；别拿 45/55 硬套；抠字翻白皮书/FAQ。",
        ],
    ),
    (
        "进展与路线图（定性）",
        [
            "形态：公示→试点→扩面（定性；无日期承诺）",
            "产品：主路径能力迭代与稳定性（定性）",
            "链上：按国/按环境逐步打开；对照已发布公告与环境标签",
            "市场：渠道试点与披露节奏对齐（试点名以已公告为准，可复读）",
            "合规与运营节奏：季度量化与内部里程碑在 NDA",
        ],
        [
            "四条并进：产品稳、链上按国开、市场跑、合规拆开。",
            "四轨为方向叙事与披露节奏，不构成业绩承诺或排期保证。",
            "数字须可溯源（来源、日期、环境标签）；口头数字与财务材料一致；条款以白皮书为准。",
        ],
    ),
    (
        "Go-to-Market",
        [
            "供给：KYB+履约+培训",
            "需求：披露口径+渠道试点",
            "渠道：把合作摊开说 · 订单资金状态分开算",
            "切入：向导/DMC 区域试点先行（定性；披露节奏随已发布公告）；阶段赢=试点与披露同轨可审计",
        ],
        [
            "渠道试点须公示；渠道分区域试点；四线叙事保持一致；对外引用预先核对。",
            "本页为打法与渠道定性，不构成收入或份额承诺。",
        ],
    ),
    (
        "竞争矩阵（定性）",
        [
            "Web2：争议弱\n典型：对账/争议难审计、约束弱",
            "链上：缺履约栈\n典型：缺旅行订单态+托管一体",
            "TravelTrust：履约栈\n典型：证据+托管+争议一体",
            "模仿成本：披露·对账·客服·嵌入切换—见 §10.2；非一夜可抄。\n「垂直 SaaS 拼盘？」：难在同一旅行订单 ID 下串齐证据→托管→争议→信誉；口径见同包 FAQ 30。\n披露：以白皮书已发布口径为准；细分对照在 NDA 沟通范围",
        ],
        [
            "费用/分配等对外表述须与同版本白皮书一致；口头对比以白皮书 + 保密材料为准。",
            "真要细比：拿白皮书和保密里那份表说话。",
            "IR 备忘·合伙人会深问顺序：FAQ 28→3→23→29→30→本页→18→20→31→16（对外 FAQ PDF 无导航行，仅 Notes 使用）。",
        ],
    ),
    (
        "典型旅程",
        [
            "Alex↔Sam（化名）：里程碑与报价→入托管",
            "履约证明与释放/裁决分支",
            "非实时清结算·45/55 脚注（分母见白皮书）",
            "教学案例·非单一真实订单",
        ],
        [
            "细节见白皮书与保密材料。",
        ],
    ),
    (
        "诉求与下一步怎么聊",
        [
            "诉求：信任结构 · 履约 · 生态",
            "披露与节奏：以已发布公告与环境标签为准",
            "下一步：平台路径 vs 基建嵌入可并行对齐；定量与法律文本走 NDA",
            "会后阅读顺序：FAQ → 白皮书（第一层披露勿与链上操作费/罚没口头混谈）",
            "「DD 增量」索引（NDA 内）：06 对照表 + 环境矩阵 + 渠道试点清单 — 仅索引，不给未披露数字。",
        ],
        [
            "常见讨论：平台与基础设施路径并存；产品侧重可嵌入履约与争议能力。",
            "治理信号与订单旅资分轨；链上与运行环境细节在保密沟通范围内展开。",
            "非实时清结算：45/55 与罚没/仲裁费用分母不同；口头叙述避免合并。",
            "定量与法律文本在保密协议后提供；以签署文本为准。",
        ],
    ),
]

SLIDES_EN: list[SlideSpec] = [
    (
        "TravelTrust",
        [
            "Trip-order fintech",
            "Custody · fulfillment · disputes for each trip order",
            "Disclosures follow published releases and the Whitepaper (product/legal anchor)",
            "Misread firewall: per-order trust + custody/disputes—not trip-pay plumbing, not mass-market OTA, not a token raise slide (same framing as Whitepaper §06 + FAQ).",
        ],
        [],
    ),
    (
        "TravelTrust",
        [
            "Trip-order fintech: align delivery and disputes around each trip order.",
            "Boundary: trip principal vs platform recurring revenue are disclosed separately; figures and terms follow the Whitepaper + NDA materials.",
            "Why it matters (qualitative): fragmented distribution needs auditable order funds + fulfillment/dispute states—embeddable rails; no revenue forecast on these slides.",
            "Why us anchor: trip funds / take / governance read on the same Whitepaper + FAQ sentences—not a generic pay-facade story.",
        ],
        [
            "Positioning (notes only): consumer UX plus embeddable custody/dispute rails—no “flip incumbents overnight” claim; aligns with slide 9 summary in notes.",
        ],
    ),
    (
        "Pain",
        [
            "Facts don’t line up",
            "Reconciliation drag; weak audit trail on disputes",
            "Traveler vs platform order funds—who may release, unclear?",
        ],
        [],
    ),
    (
        "Trip-order: align delivery and disputes · customer journey / fulfillment loop (one trip; no fee unpack here)",
        [],
        [
            "Opening line: one screen = one trip’s fulfillment path; no pools or ratios on this slide.",
            "Diagram-only trip loop; no figures beyond what’s in the pack.",
            "Vectors are a mental model only; open with environment labels: demo / testnet / prod per published releases (same discipline as slide 7).",
            "Sequence: find → book → align delivery → resolve issues → follow published disclosures.",
            "Progress and reach: stay on published disclosures; don’t imply undeclared scope or wider rollout.",
        ],
    ),
    (
        "Protocol stack · control plane — order rules & responsibility slices (same trip; rules/disclosure/env; not a second journey slide)",
        [],
        [
            "Opening line: same trip—this slide is the control stack, not a second storyline.",
            "Same single-trip logic; vertical stack shows responsibility slices (control plane); no new metrics off-pack.",
            "If the stack shows abbreviations: spell environment labels first, then layers—no implied full production before formal release.",
            "Same five-step sequence: find → book → align delivery → resolve issues → follow published disclosures.",
            "Progress and reach: stay on published disclosures; don’t imply undeclared scope. Names and splits—after the meeting; avoid blanket live claims before formal public release.",
        ],
    ),
    (
        "Why now (qualitative)",
        [
            "Demand: be real, not loud. (structure: Whitepaper §2.4)",
            "Supply: reputation travels; align releases to published communications + labels",
            "Regulatory: don’t mix traveler order funds state with platform earnings. (Whitepaper)",
            "Window: splintered channels—align delivery first, then distribution. (same section) Replay anchor: published §06 refs + release dates + environment labels.",
        ],
        [
            "Timing replay: window narrative reads against published §06 section refs, release dates, and env labels; industry ratio tables stay in NDA.",
            "Industry comps and tables usually need NDA; this isn’t a “buy now” nudge.",
        ],
    ),
    (
        "Three lines to remember",
        [
            "Platform take is ours; order funds state isn’t.",
            "Governance signal: votes, not trip pay.",
            "Disclosure cadence: published releases + version/env labels (details: Whitepaper). · Moat leans on replayable disclosure + recon.",
        ],
        [
            "Line 1: recurring revenue from service take; trip principal is not our own income.",
            "Lines 2–3 as on slide; don’t unpack disclosure rails on this slide.",
            "Three sentences carry the point; detail after the meeting.",
        ],
    ),
    (
        "One order: delivery aligned?",
        [
            "One order: book through to aligned delivery?",
            "Scope: follow published releases and environment labels.",
            "Materials: FAQ + Whitepaper; keep published 45/55 allocatable-fee split verbally separate from chain ops fees / penalties.",
        ],
        [
            "Walk one booking to aligned delivery; don’t go past published releases and environment labels.",
        ],
    ),
    (
        "System strengths (qualitative view: embeddable trust stack; moat in disclosure & reconciliation)",
        [
            "Fulfillment · order alignment · delivery & disputes · clear resolution · published releases/labels",
        ],
        [
            "Summary (notes only): consumer UX + embeddable custody/dispute rails—no overnight flip claim.",
            "On-screen pillars only; keep one visual row for a fast scan.",
        ],
    ),
    (
        "Published path · disclosure boundary (outer-ring disclosure | FeeRouter — published fee routing | on-slide published path only)",
        [],
        [
            "Opening line: published disclosure boundary (outer ring); align denominators—not operating P&L.",
            "Disclosure boundary diagram; slide stays 100/45/55 only.",
            "Verbal tripod (same as FAQ / Whitepaper): published 45/55 is the first split of allocatable order fees; gas/arbitration/slashing use different denominators.",
            "Disclosure·45%: country bucket; inner split in signed addenda.",
            "Disclosure·55%: Global Pool; inner split in signed addendum.",
            "Gas/arbitration isn’t this bucket; don’t mash with 45/55; wording in Whitepaper + FAQ.",
        ],
    ),
    (
        "Roadmap & timing (qualitative)",
        [
            "Shape: disclose → pilot → expand (qualitative; no dated promise)",
            "Product: core path capability iteration (qualitative)",
            "On-chain: roll out by country/environment—published releases + environment labels",
            "GTM: channel pilots aligned to disclosure cadence (pilot names only as already published—replayable)",
            "Compliance & ops cadence: quarterly metrics and internal milestones under NDA",
        ],
        [
            "Four tracks: ship product, turn on-chain by country, grow channels, split compliance.",
            "Tracks are directional narrative—not performance guarantees or dated commitments.",
            "Numbers need source, date, and whether it’s demo, test, or live; check verbal figures against Finance handouts; governance signal and disclosure splits are spelled out in the Whitepaper.",
        ],
    ),
    (
        "Go-to-market",
        [
            "Supply: KYB + fulfillment + training",
            "Demand: disclosure + channel pilots",
            "Channels: spell out channel pilots · order funds state separate",
            "Wedge: guide/DMC regional pilots first (qualitative; disclosure tracks published releases + labels); stage win = pilots + disclosure cadence auditable together",
        ],
        [
            "Report channel pilots; DMC pilots; keep the four-track story coherent; quotes pre-cleared.",
            "GTM here is qualitative positioning—not a revenue or share commitment.",
        ],
    ),
    (
        "Competitive matrix (qualitative)",
        [
            "Web2: weak disputes\nPattern: recon pain; weak auditable disputes",
            "On-chain: no fulfillment stack\nPattern: travel order state + custody rarely integrated",
            "TravelTrust: fulfillment stack\nPattern: evidence + escrow + disputes in one integrated trip-order flow",
            "Switching cost: disclosure, recon, support, embed cutovers—see §10.2; not overnight.\n“Seven best-of-breed SaaS” bolt-ons: rarely one **trip-order spine** (evidence→escrow→dispute→reputation); see pack **FAQ #30**.\nDisclosure: follow the published Whitepaper; line-item tables in the NDA room",
        ],
        [
            "Fees/allocations must match the same Whitepaper release; comparisons use Whitepaper + confidential tables.",
            "If you want the receipts: Whitepaper + the NDA room tables.",
            "IR memo · partner-meeting deep-dive order: FAQ 28→3→23→29→30→this slide→18→20→31→16 (not printed in the public FAQ PDF—notes only).",
        ],
    ),
    (
        "Narrative journey",
        [
            "Alex↔Sam (pseudonyms): milestones & quote → escrow visible",
            "Fulfillment proofs and release / ruling branch",
            "Non-real-time settlement footnote (denominators: Whitepaper)",
            "Teaching case—not one live order",
        ],
        [
            "Details live in the Whitepaper and NDA room.",
        ],
    ),
    (
        "Asks & next conversation",
        [
            "Asks: trust · fulfillment · ecosystem",
            "Disclosure & cadence: published releases + environment labels",
            "Next: platform vs rails can be discussed in parallel; quant + legal text after NDA",
            "After the meeting: FAQ → Whitepaper (keep published 45/55 split verbally separate from chain ops fees / penalties)",
            "“DD delta” index (under NDA): §06 cross-walk + env matrix + channel pilot list—index only; no undisclosed figures on this slide.",
        ],
        [
            "Typical question: platform vs rails—both can fit; UX plus embeddable fulfillment/disputes.",
            "Governance signal and trip order funds state stay separate; chain details and environment under NDA.",
            "Non-real-time settlement: 45/55 is not the same denominator as penalties/arbitration charges—don’t blend them aloud.",
            "Quant and legal text after NDA; governed by signed documents only.",
        ],
    ),
]


def _release() -> str:
    data = json.loads(ANCHORS.read_text(encoding="utf-8"))
    r = data.get("release")
    if not r:
        raise SystemExit("numeric anchors JSON missing release")
    return str(r)


def _try_soffice(pptx: Path, out_dir: Path) -> bool:
    exe_names = ("soffice", "soffice.exe")
    candidates: list[Path] = []
    pf = os.environ.get("PROGRAMFILES", r"C:\Program Files")
    pf86 = os.environ.get("PROGRAMFILES(X86)", r"C:\Program Files (x86)")
    for base in (Path(pf) / "LibreOffice" / "program", Path(pf86) / "LibreOffice" / "program"):
        candidates.append(base / "soffice.exe")

    for c in candidates:
        if c.is_file():
            cmd = [str(c), "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx)]
            try:
                subprocess.run(cmd, check=True, capture_output=True, timeout=120)
                return True
            except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError):
                pass

    for name in exe_names:
        cmd = [name, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx)]
        try:
            subprocess.run(cmd, check=True, capture_output=True, timeout=120)
            return True
        except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired, OSError):
            continue
    return False


def _build_one(lang: str, slides: list[SlideSpec], release: str) -> Path:
    try:
        from pptx import Presentation  # noqa: PLC0415
        from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE  # noqa: PLC0415
        from pptx.util import Inches, Pt  # noqa: PLC0415
    except ImportError as e:
        raise SystemExit("pip install python-pptx") from e

    tk = ids.active_theme()
    InchesF = Inches

    _tools = Path(__file__).resolve().parent
    if str(_tools) not in sys.path:
        sys.path.insert(0, str(_tools))
    import ir_brand_vector_charts as irv  # noqa: PLC0415

    prs = Presentation()
    prs.slide_width = Inches(ids.SLIDE_W_IN)
    prs.slide_height = Inches(ids.SLIDE_H_IN)

    blank = prs.slide_layouts[6]
    total = len(slides)
    cn = lang == "CN"
    foot_short = (
        f"TravelTrust · Release {release} · General information only — not an offer"
        if not cn
        else f"TravelTrust · {release} · 一般性信息 — 非要约"
    )

    for idx, spec in enumerate(slides, start=1):
        title, bullets, note_lines = _unpack_slide(spec)
        deck_detail = "\n".join(note_lines) if note_lines else ""

        slide = prs.slides.add_slide(blank)
        if idx == 1:
            ids.add_cover(
                slide,
                tk,
                title=title,
                subtitle_lines=bullets,
                release=release,
                lang=lang,
                slide_idx=idx,
                total=total,
                MSO_SHAPE=MSO_SHAPE,
                InchesF=InchesF,
            )
            slide.notes_slide.notes_text_frame.text = ids_story.speaker_note_text(
                slide_idx=idx, title=title, lang=lang, deck_detail=deck_detail
            )
            ids.apply_fade_transition(slide)
            continue

        kind = _slide_layout_kind(title, bullets)
        beat = ids_story.beat_for_slide_idx(idx)
        ids.add_narrative_rail(
            slide, tk, beat.macro_arc, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF
        )
        eyebrow = ids_story.eyebrow_for_slide(idx, lang)
        fee_split = _fee_slide_chrome_split(title, lang) if kind == "vector_fee" else None
        long_split = _pitch_chrome_long_split(title, lang, idx)
        if fee_split:
            top = ids.add_chrome_bar(
                slide,
                tk,
                fee_split[0],
                eyebrow=eyebrow,
                lang=lang,
                MSO_SHAPE=MSO_SHAPE,
                InchesF=InchesF,
                title_subline=fee_split[1],
            )
        elif long_split:
            top = ids.add_chrome_bar(
                slide,
                tk,
                long_split[0],
                eyebrow=eyebrow,
                lang=lang,
                MSO_SHAPE=MSO_SHAPE,
                InchesF=InchesF,
                title_subline=long_split[1],
            )
        else:
            top = ids.add_chrome_bar(
                slide, tk, title, eyebrow=eyebrow, lang=lang, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF
            )

        pillars = _pillar_parts(title, bullets)

        if kind == "hero" and bullets:
            ids.add_hero_statement(slide, tk, bullets[0], top_y_in=top, lang=lang, InchesF=InchesF)
        elif kind == "pain_grid" and bullets:
            cols = 3 if len(bullets) == 3 else 2
            ids.add_bullet_cards_2col(
                slide, tk, bullets, top_y_in=top, lang=lang, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF, columns=cols
            )
        elif kind == "pillars" and pillars:
            ids.add_pillar_strip(slide, tk, pillars, top_y_in=top, lang=lang, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF)
        elif kind == "matrix":
            trip = _split_competitive_columns(bullets)
            if trip:
                ids.add_matrix_three_plus_note(
                    slide,
                    tk,
                    trip[0],
                    trip[1],
                    trip[2],
                    trip[3],
                    top_y_in=top,
                    lang=lang,
                    MSO_SHAPE=MSO_SHAPE,
                    InchesF=InchesF,
                )
            else:
                ids.add_body_bullets(slide, tk, bullets, top_y_in=top, lang=lang, InchesF=InchesF)
        elif kind == "faq" and bullets:
            ids.add_faq_row_cards(slide, tk, bullets, top_y_in=top, lang=lang, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF)
        elif kind == "roadmap" and bullets:
            ids.add_roadmap_tracks(slide, tk, bullets, top_y_in=top, lang=lang, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF)
        elif kind == "cards_2col" and bullets:
            ids.add_bullet_cards_2col(slide, tk, bullets, top_y_in=top, lang=lang, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF, columns=2)
        elif kind == "vector_loop":
            irv.pptx_add_travel_loop(slide, lang, MSO_SHAPE, MSO_CONNECTOR, InchesF, Pt, tk.accent, tk.ink, tk.surface)
        elif kind == "vector_stack":
            irv.pptx_add_protocol_stack(slide, lang, MSO_SHAPE, InchesF, Pt, tk.bg_deep, tk.accent, tk.ink, tk.surface)
        elif kind == "vector_fee":
            irv.pptx_add_fee_router_layer1(slide, lang, MSO_SHAPE, MSO_CONNECTOR, InchesF, Pt, tk.accent, tk.ink, tk.surface)
        elif bullets:
            ids.add_body_bullets(slide, tk, bullets, top_y_in=top, lang=lang, InchesF=InchesF)

        ids.add_footer(slide, tk, foot_short, lang=lang, slide_idx=idx, total=total, InchesF=InchesF)
        slide.notes_slide.notes_text_frame.text = ids_story.speaker_note_text(
            slide_idx=idx, title=title, lang=lang, deck_detail=deck_detail
        )
        ids.apply_fade_transition(slide)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    edit = ihr.maintainer_deck_editable_dir()
    edit.mkdir(parents=True, exist_ok=True)
    suffix = "CN" if lang == "CN" else "EN"
    path = edit / f"04-PitchDeck-v{release}-{suffix}.pptx"
    prs.save(str(path))
    return path


def _pdf_via_fpdf(lang: str, slides: list[SlideSpec], release: str, out_pdf: Path) -> None:
    """Fallback PDF when LibreOffice/PowerPoint export is unavailable (same copy as PPTX)."""
    try:
        from fpdf import FPDF  # noqa: PLC0415
    except ImportError as e:
        raise SystemExit("pip install fpdf2 (required for PDF fallback)") from e

    simhei = Path(r"C:/Windows/Fonts/simhei.ttf")
    arial = Path(r"C:/Windows/Fonts/arial.ttf")
    use_cn_font = lang == "CN" and simhei.is_file()

    class _P(FPDF):
        pass

    pdf = _P()
    pdf.set_auto_page_break(auto=True, margin=14)
    if use_cn_font:
        pdf.add_font("TT", "", str(simhei))
        body_font = "TT"
    elif arial.is_file():
        pdf.add_font("EN", "", str(arial))
        body_font = "EN"
    else:
        body_font = "Helvetica"
    footer = f"TravelTrust · Release {release} · General information only - not an offer"

    _tools = Path(__file__).resolve().parent
    if str(_tools) not in sys.path:
        sys.path.insert(0, str(_tools))
    import ir_brand_vector_charts as irv  # noqa: PLC0415

    for idx, spec in enumerate(slides, start=1):
        title, bullets, _ = _unpack_slide(spec)
        pdf.add_page()
        pdf.set_font(body_font, size=16 if use_cn_font else 14)
        pdf.multi_cell(0, 8, title, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(2)
        pdf.set_font(body_font, size=11)
        y0 = pdf.get_y()
        if not bullets:
            if (
                "单笔订单对齐交付与争议" in title
                or "trip-order: align delivery and disputes" in title.lower()
                or "主链" in title
                or "deal loop" in title.lower()
            ):
                irv.fpdf_travel_loop(pdf, body_font, lang, y0)
            elif (
                "订单与规则分几层" in title
                or "钱和规则分几层" in title
                or "trip-order & rules stack" in title.lower()
                or "money & rules stack" in title.lower()
                or "协议栈" in title
                or "Protocol stack" in title
                or "order rules |" in title.lower()
            ):
                irv.fpdf_protocol_stack(pdf, body_font, lang, y0)
            elif (
                "feerouter" in title.lower()
                or "协议经济·第一层" in title
                or ("protocol economics" in title.lower() and "layer 1" in title.lower())
                or "disclosure boundary" in title.lower()
                or "披露边界" in title
            ):
                irv.fpdf_fee_router_layer1(pdf, body_font, lang, y0)
        else:
            for b in bullets:
                line = b.replace("→", "->") if lang == "EN" else b
                line = line.replace("\u2014", "-").replace("\u2013", "-")
                pdf.multi_cell(0, 6, line, new_x="LMARGIN", new_y="NEXT")
                pdf.ln(1)
        pdf.set_y(-12)
        pdf.set_font(body_font, size=8)
        pdf.set_text_color(120, 120, 120)
        pdf.cell(0, 6, footer, new_x="RIGHT", new_y="TOP")
        pdf.set_text_color(0, 0, 0)

    pdf.output(str(out_pdf))


def main() -> int:
    release = _release()
    ids_story.assert_story_aligned(len(SLIDES_CN))
    if len(SLIDES_EN) != len(SLIDES_CN):
        raise SystemExit("SLIDES_EN length must match SLIDES_CN")
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ihr.purge_traveltrust_legacy_binaries(OUT_DIR)
    ihr.set_fundraising_root(ROOT / "docs" / "fundraising")
    ihr.purge_pitch_deck_for_release(OUT_DIR, release)
    cn = _build_one("CN", SLIDES_CN, release)
    en = _build_one("EN", SLIDES_EN, release)
    print(f"OK: {cn}")
    print(f"OK: {en}")

    soffice = find_soffice()
    for p in (cn, en):
        pdf_path = p.with_suffix(".pdf")
        convert_pptx_to_pdf(soffice, p, OUT_DIR)
        if not pdf_path.is_file():
            raise SystemExit(f"expected Pitch Deck PDF at {pdf_path}")
        print(f"OK: PDF via LibreOffice -> {pdf_path}")
    extras = ihr.purge_export_ready_handoff_extras(OUT_DIR, release)
    if extras:
        print("OK: purged from export-ready:", ", ".join(extras))
    ihr.write_start_here(OUT_DIR, release)
    print(f"OK: {OUT_DIR / '00-START-HERE.txt'}")
    demo_dir = OUT_DIR / "demo"
    if os.environ.get("TRAVELTRUST_IR_SLIDE_AUDIT", "").strip().lower() in ("1", "true", "yes"):
        audit_path = ids_audit.write_audit_bundle(
            OUT_DIR, release, _slides_audit_pairs(SLIDES_CN), _slides_audit_pairs(SLIDES_EN)
        )
        print(f"OK: slide audit -> {audit_path}")
    else:
        print("OK: slide audit skipped (set TRAVELTRUST_IR_SLIDE_AUDIT=1 to write internal/37-IR-Pitch-Deck-Slide-Audit-RUN.md)")
    # IC appendix is internal-only (deck-editable); not copied to export-ready.
    if os.environ.get("TRAVELTRUST_BUILD_IC_INTERNAL", "").strip().lower() in ("1", "true", "yes"):
        ic_py = ROOT / "scripts" / "tools" / "build-investor-ic-appendix.py"
        if ic_py.is_file():
            r = subprocess.run([sys.executable, str(ic_py)], cwd=str(ROOT), capture_output=True, text=True)
            if r.stdout:
                print(r.stdout, end="")
            if r.returncode != 0:
                raise SystemExit(r.stderr or "IC appendix build failed")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
