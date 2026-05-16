#!/usr/bin/env python3
"""
Build TravelTrust **IC / Partner Summary** appendix PPTX (+ PDF) — **8 slides** (cover + 7 IC themes).

Does **not** modify the main 15-slide Pitch Deck. Same design tokens as `ir_deck_design_system.py`.
No new economic numbers (references to Whitepaper layer-1 100%/45/55% only where listed).

Outputs `internal/deck-editable/04-IC-Memo-v{release}-CN|EN.pptx` (+ optional PDF there). **Not** shipped in export-ready.

When `TRAVELTRUST_SKIP_IC=1`, `build-investor-pitch-deck.py` skips invoking this script; run this file directly anytime.
"""
from __future__ import annotations

import json
import os
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
ANCHORS = ROOT / "registry" / "fundraising-external-numeric-anchors.v1.json"
OUT_DIR = ROOT / "docs" / "fundraising" / "internal" / "deck-editable"

_TOOLS = Path(__file__).resolve().parent
if str(_TOOLS) not in sys.path:
    sys.path.insert(0, str(_TOOLS))
import investor_handoff_layout as ihr  # noqa: E402
import ir_deck_design_system as ids  # noqa: E402
from ir_ic_appendix import (  # noqa: E402
    IC_EYEBROW_CN,
    IC_EYEBROW_EN,
    IC_RAIL_MACROS,
    IC_SLIDES_CN,
    IC_SLIDES_EN,
    SlideSpec,
    ic_speaker_note,
)
from investor_ir_office import convert_pptx_to_pdf, find_soffice  # noqa: E402


def _unpack_slide(spec: SlideSpec) -> tuple[str, list[str], list[str]]:
    title, vis, notes = spec
    return str(title), list(vis), list(notes)


def _split_competitive_columns(bullets: list[str]) -> tuple[tuple[str, str], tuple[str, str], tuple[str, str], str] | None:
    if len(bullets) < 4:
        return None

    def one(line: str) -> tuple[str, str]:
        for sep in ("：", ":"):
            if sep in line:
                a, b = line.split(sep, 1)
                return a.strip(), b.strip()
        return line.strip(), ""

    return one(bullets[0]), one(bullets[1]), one(bullets[2]), bullets[3]


def _pillar_parts(title: str, bullets: list[str]) -> list[str] | None:
    if ("优势" in title or "moat" in title.lower()) and len(bullets) == 1 and "·" in bullets[0]:
        s = bullets[0].replace(" · ", "·").replace(" • ", "·")
        parts = [p.strip() for p in s.split("·") if p.strip()]
        return parts or None
    return None


def _ic_slide_layout_kind(title: str, bullets: list[str]) -> str:
    t = title.lower()
    if "竞争" in title or "competition" in t:
        return "matrix"
    if "dd" in t or "尽调" in title:
        return "faq"
    if len(bullets) >= 4:
        return "cards_2col"
    if len(bullets) >= 3:
        return "cards_2col"
    if _pillar_parts(title, bullets):
        return "pillars"
    return "linear"


def _release() -> str:
    data = json.loads(ANCHORS.read_text(encoding="utf-8"))
    r = data.get("release")
    if not r:
        raise SystemExit("numeric anchors JSON missing release")
    return str(r)


def _build_ic_one(lang: str, slides: list[SlideSpec], release: str) -> Path:
    try:
        from pptx import Presentation  # noqa: PLC0415
        from pptx.enum.shapes import MSO_SHAPE  # noqa: PLC0415
        from pptx.util import Inches, Pt  # noqa: PLC0415
    except ImportError as e:
        raise SystemExit("pip install python-pptx") from e

    tk = ids.active_theme()
    InchesF = Inches
    total = len(slides)
    cn = lang == "CN"
    foot_short = (
        f"TravelTrust · IC appendix · Release {release} · General information only — not an offer"
        if not cn
        else f"TravelTrust · IC 附录 · {release} · 一般性信息 — 非要约"
    )

    prs = Presentation()
    prs.slide_width = Inches(ids.SLIDE_W_IN)
    prs.slide_height = Inches(ids.SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    for idx, spec in enumerate(slides, start=1):
        title, bullets, note_lines = _unpack_slide(spec)
        deck_detail = "\n".join(note_lines) if note_lines else ""
        slide = prs.slides.add_slide(blank)

        if idx == 1:
            ids.add_cover(
                slide,
                tk,
                title=title,
                subtitle_lines=bullets,
                release=release,
                lang=lang,
                slide_idx=idx,
                total=total,
                MSO_SHAPE=MSO_SHAPE,
                InchesF=InchesF,
            )
            slide.notes_slide.notes_text_frame.text = ic_speaker_note(
                slide_idx=idx, title=title, lang=lang, deck_detail=deck_detail
            )
            ids.apply_fade_transition(slide)
            continue

        macro = IC_RAIL_MACROS[idx - 2]
        ids.add_narrative_rail(slide, tk, macro, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF)
        eyebrow = IC_EYEBROW_CN[idx - 2] if cn else IC_EYEBROW_EN[idx - 2]
        top = ids.add_chrome_bar(slide, tk, title, eyebrow=eyebrow, lang=lang, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF)

        kind = _ic_slide_layout_kind(title, bullets)
        pillars = _pillar_parts(title, bullets)

        if kind == "pillars" and pillars:
            ids.add_pillar_strip(slide, tk, pillars, top_y_in=top, lang=lang, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF)
        elif kind == "matrix":
            trip = _split_competitive_columns(bullets)
            if trip:
                ids.add_matrix_three_plus_note(
                    slide,
                    tk,
                    trip[0],
                    trip[1],
                    trip[2],
                    trip[3],
                    top_y_in=top,
                    lang=lang,
                    MSO_SHAPE=MSO_SHAPE,
                    InchesF=InchesF,
                )
            else:
                ids.add_body_bullets(slide, tk, bullets, top_y_in=top, lang=lang, InchesF=InchesF)
        elif kind == "faq" and bullets:
            ids.add_faq_row_cards(slide, tk, bullets, top_y_in=top, lang=lang, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF)
        elif kind == "cards_2col" and bullets:
            ids.add_bullet_cards_2col(slide, tk, bullets, top_y_in=top, lang=lang, MSO_SHAPE=MSO_SHAPE, InchesF=InchesF, columns=2)
        elif bullets:
            ids.add_body_bullets(slide, tk, bullets, top_y_in=top, lang=lang, InchesF=InchesF)

        ids.add_footer(slide, tk, foot_short, lang=lang, slide_idx=idx, total=total, InchesF=InchesF)
        slide.notes_slide.notes_text_frame.text = ic_speaker_note(
            slide_idx=idx, title=title, lang=lang, deck_detail=deck_detail
        )
        ids.apply_fade_transition(slide)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    suffix = "CN" if lang == "CN" else "EN"
    path = OUT_DIR / f"04-IC-Memo-v{release}-{suffix}.pptx"
    prs.save(str(path))
    return path


def main() -> int:
    release = _release()
    if len(IC_SLIDES_CN) != len(IC_SLIDES_EN):
        raise SystemExit("IC_SLIDES_CN and IC_SLIDES_EN must have the same length")
    if len(IC_RAIL_MACROS) != len(IC_SLIDES_CN) - 1:
        raise SystemExit("IC_RAIL_MACROS must have len(IC slides) - 1 (cover has no rail row)")
    if len(IC_EYEBROW_CN) != len(IC_RAIL_MACROS) or len(IC_EYEBROW_EN) != len(IC_RAIL_MACROS):
        raise SystemExit("IC eyebrows must match rail macro count")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ihr.set_fundraising_root(ROOT / "docs" / "fundraising")
    export_ready = ROOT / "docs" / "fundraising" / "external" / "export-ready"
    ihr.purge_ic_memo_for_release(export_ready, release)
    ihr.purge_ic_memo_for_release(OUT_DIR, release)
    cn = _build_ic_one("CN", IC_SLIDES_CN, release)
    en = _build_ic_one("EN", IC_SLIDES_EN, release)
    print(f"OK: {cn}")
    print(f"OK: {en}")

    soffice = find_soffice()
    for p in (cn, en):
        pdf_path = p.with_suffix(".pdf")
        convert_pptx_to_pdf(soffice, p, OUT_DIR)
        if not pdf_path.is_file():
            raise SystemExit(f"expected IC Memo PDF at {pdf_path}")
        print(f"OK: PDF via LibreOffice -> {pdf_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
